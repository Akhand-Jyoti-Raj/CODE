from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title Slide
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "Desalination of Water"
subtitle.text = "Addressing Water Scarcity Through Innovation"

# Slide 2: Introduction
slide_intro = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_intro.shapes.title, slide_intro.placeholders[1]
title.text = "Introduction"
content.text = "Overview of the global water crisis\nImportance of desalination in addressing water scarcity"

# Slide 3: Water Scarcity
slide_scarcity = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_scarcity.shapes.title, slide_scarcity.placeholders[1]
title.text = "Water Scarcity"
content.text = "Statistics on global water scarcity\nImpact on communities and ecosystems"

# Slide 4: What is Desalination?
slide_desalination = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_desalination.shapes.title, slide_desalination.placeholders[1]
title.text = "What is Desalination?"
content.text = "Definition of desalination\nImportance of converting seawater to freshwater"

# Slide 5: Desalination Technologies
slide_technologies = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_technologies.shapes.title, slide_technologies.placeholders[1]
title.text = "Desalination Technologies"
content.text = "Overview of common desalination methods:\n- Reverse Osmosis\n- Multi-Stage Flash Distillation\n- Multi-Effect Distillation\n- Electrodialysis"

# Add more slides similarly for each section of your presentation...

# Save the presentation
prs.save("Desalination_Presentation.pptx")
